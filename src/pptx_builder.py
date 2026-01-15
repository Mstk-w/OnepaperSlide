"""
OnePaperSlide PPTX生成モジュール
レイアウトデータからPowerPointファイルを生成
"""

from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
from typing import Dict, Any
from custom_types import SlideLayoutData

from logging_config import get_logger

# モジュールロガー
logger = get_logger("pptx_builder")


class Colors:
    """デザイン仕様に基づくカラーパレット"""
    PRIMARY = RGBColor(0x2B, 0x6C, 0xB0)      # 青 - 見出し・アクセント
    SECONDARY = RGBColor(0x4A, 0x55, 0x68)    # グレー - 本文
    BACKGROUND = RGBColor(0xFF, 0xFF, 0xFF)   # 白
    ACCENT_BG = RGBColor(0xEB, 0xF8, 0xFF)    # 薄い青 - ボックス背景
    BORDER = RGBColor(0xE2, 0xE8, 0xF0)       # 薄いグレー - 境界線
    ALERT = RGBColor(0xC5, 0x30, 0x30)        # 赤 - 警告

    @staticmethod
    def from_hex(hex_color: str) -> RGBColor:
        """16進数文字列からRGBColorを生成"""
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )


class Fonts:
    """タイポグラフィ設定"""
    FAMILY = "メイリオ"
    TITLE = 28
    SECTION_HEADER = 18
    BODY = 14
    FOOTER = 10
    TABLE_HEADER = 12
    TABLE_BODY = 11


def add_text_box(
    slide,
    text: str,
    x_mm: float,
    y_mm: float,
    width_mm: float,
    height_mm: float,
    font_size_pt: int = Fonts.BODY,
    font_bold: bool = False,
    font_color: RGBColor = None,
    alignment: PP_ALIGN = PP_ALIGN.LEFT
):
    """テキストボックスを追加"""
    if font_color is None:
        font_color = Colors.SECONDARY

    shape = slide.shapes.add_textbox(
        Mm(x_mm), Mm(y_mm),
        Mm(width_mm), Mm(height_mm)
    )

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = Fonts.FAMILY
    p.font.size = Pt(font_size_pt)
    p.font.bold = font_bold
    p.font.color.rgb = font_color
    p.alignment = alignment

    return shape


def add_section_header(
    slide,
    text: str,
    x_mm: float,
    y_mm: float,
    width_mm: float
):
    """セクション見出しを追加（アクセントバー付き）"""
    # アクセントバー（左端の青い縦線）
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Mm(x_mm), Mm(y_mm),
        Mm(3), Mm(8)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = Colors.PRIMARY
    bar.line.fill.background()

    # 見出しテキスト
    add_text_box(
        slide,
        text,
        x_mm + 5,
        y_mm,
        width_mm - 5,
        10,
        font_size_pt=Fonts.SECTION_HEADER,
        font_bold=True,
        font_color=Colors.PRIMARY
    )


def add_bullet_list(
    slide,
    items: list,
    x_mm: float,
    y_mm: float,
    width_mm: float,
    height_mm: float
):
    """箇条書きリストを追加"""
    shape = slide.shapes.add_textbox(
        Mm(x_mm), Mm(y_mm),
        Mm(width_mm), Mm(height_mm)
    )

    tf = shape.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = item if isinstance(item, str) else item.get("text", "")
        p.text = "・" + text
        p.font.name = Fonts.FAMILY
        p.font.size = Pt(Fonts.BODY)
        p.font.color.rgb = Colors.SECONDARY
        p.space_after = Pt(4)


def add_box_with_background(
    slide,
    x_mm: float,
    y_mm: float,
    width_mm: float,
    height_mm: float,
    fill_color: RGBColor = None,
    border_color: RGBColor = None
):
    """角丸矩形ボックスを追加"""
    if fill_color is None:
        fill_color = Colors.ACCENT_BG
    if border_color is None:
        border_color = Colors.BORDER

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Mm(x_mm), Mm(y_mm),
        Mm(width_mm), Mm(height_mm)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.5)

    return shape


def add_table(
    slide,
    x_mm: float,
    y_mm: float,
    width_mm: float,
    columns: list,
    rows: list
):
    """表を追加"""
    if not columns:
        return None

    row_count = len(rows) + 1
    col_count = len(columns)

    table = slide.shapes.add_table(
        row_count, col_count,
        Mm(x_mm), Mm(y_mm),
        Mm(width_mm), Mm(row_count * 8)
    ).table

    col_width = Mm(width_mm / col_count)
    for col in table.columns:
        col.width = col_width

    # ヘッダー行
    for i, col_name in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = Colors.PRIMARY

        p = cell.text_frame.paragraphs[0]
        p.font.name = Fonts.FAMILY
        p.font.size = Pt(Fonts.TABLE_HEADER)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # データ行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_value in enumerate(row_data):
            if col_idx >= col_count:
                continue
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_value)

            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = Colors.ACCENT_BG
            else:
                cell.fill.background()

            p = cell.text_frame.paragraphs[0]
            p.font.name = Fonts.FAMILY
            p.font.size = Pt(Fonts.TABLE_BODY)
            p.font.color.rgb = Colors.SECONDARY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def add_triangle_arrow(
    slide,
    x_mm: float,
    y_mm: float,
    direction: str = "right"
):
    """三角形の矢印を追加"""
    rotation = {"right": 90, "down": 180, "left": 270, "up": 0}[direction]

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Mm(x_mm), Mm(y_mm),
        Mm(8), Mm(8)
    )
    shape.rotation = rotation
    shape.fill.solid()
    shape.fill.fore_color.rgb = Colors.PRIMARY
    shape.line.fill.background()

    return shape


def add_flowchart(
    slide,
    steps: list,
    x_mm: float,
    y_mm: float,
    width_mm: float,
    direction: str = "h"
):
    """フローチャートを追加"""
    step_count = len(steps)
    if step_count == 0:
        return

    if direction == "h":
        step_width = (width_mm - (step_count - 1) * 15) / step_count
        step_height = 25

        for i, step in enumerate(steps):
            step_x = x_mm + i * (step_width + 15)
            text = step.get("text", step) if isinstance(step, dict) else step

            add_box_with_background(slide, step_x, y_mm, step_width, step_height)
            add_text_box(
                slide, text,
                step_x + 2, y_mm + 5,
                step_width - 4, step_height - 10,
                font_size_pt=12,
                alignment=PP_ALIGN.CENTER
            )

            if i < step_count - 1:
                add_triangle_arrow(
                    slide,
                    step_x + step_width + 3,
                    y_mm + step_height / 2 - 4,
                    "right"
                )
    else:
        step_width = width_mm - 20
        step_height = 18

        for i, step in enumerate(steps):
            step_y = y_mm + i * (step_height + 12)
            text = step.get("text", step) if isinstance(step, dict) else step

            add_box_with_background(slide, x_mm + 10, step_y, step_width, step_height)
            add_text_box(
                slide, text,
                x_mm + 12, step_y + 2,
                step_width - 4, step_height - 4,
                font_size_pt=12,
                alignment=PP_ALIGN.CENTER
            )

            if i < step_count - 1:
                add_triangle_arrow(
                    slide,
                    x_mm + 10 + step_width / 2 - 4,
                    step_y + step_height + 2,
                    "down"
                )


def auto_shrink_text(
    text: str,
    box_width_mm: float,
    box_height_mm: float,
    initial_font_size_pt: int,
    min_font_size_pt: int = 8
) -> int:
    """テキストがボックスに収まるフォントサイズを計算"""
    font_size = initial_font_size_pt

    while font_size >= min_font_size_pt:
        char_width_mm = font_size * 0.35
        chars_per_line = int(box_width_mm / char_width_mm) if char_width_mm > 0 else 1
        text_length = len(text)
        required_lines = (text_length + chars_per_line - 1) // chars_per_line
        line_height_mm = font_size * 0.42
        required_height = required_lines * line_height_mm

        if required_height <= box_height_mm:
            return font_size

        font_size -= 1

    return min_font_size_pt


def add_text_box_with_auto_shrink(
    slide,
    text: str,
    x_mm: float,
    y_mm: float,
    width_mm: float,
    height_mm: float,
    initial_font_size_pt: int = Fonts.BODY,
    min_font_size_pt: int = 8,
    **kwargs
):
    """Auto-Shrink付きテキストボックス"""
    font_size = auto_shrink_text(
        text, width_mm, height_mm,
        initial_font_size_pt, min_font_size_pt
    )

    return add_text_box(
        slide, text,
        x_mm, y_mm, width_mm, height_mm,
        font_size_pt=font_size,
        **kwargs
    )


def render_header(slide, header_data: dict):
    """ヘッダーを描画"""
    title = header_data.get("title", {})
    add_text_box(
        slide,
        title.get("text", ""),
        title.get("x_mm", 15),
        title.get("y_mm", 10),
        title.get("width_mm", 390),
        title.get("height_mm", 15),
        font_size_pt=title.get("font_size_pt", Fonts.TITLE),
        font_bold=title.get("font_bold", True),
        font_color=Colors.PRIMARY
    )

    subtitle = header_data.get("subtitle", {})
    if subtitle.get("text"):
        add_text_box(
            slide,
            subtitle["text"],
            subtitle.get("x_mm", 15),
            subtitle.get("y_mm", 28),
            subtitle.get("width_mm", 390),
            subtitle.get("height_mm", 7),
            font_size_pt=subtitle.get("font_size_pt", 14),
            font_color=Colors.SECONDARY
        )


def render_section(slide, section_data: dict):
    """セクションを描画"""
    section_type = section_data.get("type", "text_block")

    # ヘッダー描画
    header = section_data.get("header", {})
    if header.get("text"):
        add_section_header(
            slide,
            header["text"],
            header["x_mm"],
            header["y_mm"],
            header["width_mm"]
        )

    # タイプ別描画
    if section_type == "bullets":
        render_bullets(slide, section_data)
    elif section_type == "table":
        render_table(slide, section_data)
    elif section_type == "flowchart":
        render_flowchart(slide, section_data)
    elif section_type == "kpi_box":
        render_kpi_box(slide, section_data)
    else:
        render_text_block(slide, section_data)


def render_bullets(slide, data: dict):
    """箇条書きを描画"""
    items = data.get("items", [])
    if not items:
        return

    first_item = items[0]
    x_mm = first_item.get("x_mm", 20)
    y_mm = first_item.get("y_mm", 50)
    width_mm = first_item.get("width_mm", 180)

    texts = [item.get("text", item) if isinstance(item, dict) else item
             for item in items]

    add_bullet_list(
        slide, texts,
        x_mm, y_mm, width_mm,
        height_mm=len(items) * 8 + 5
    )


def render_table(slide, data: dict):
    """表を描画"""
    table_data = data.get("table", {})
    add_table(
        slide,
        table_data.get("x_mm", 20),
        table_data.get("y_mm", 50),
        table_data.get("width_mm", 180),
        table_data.get("columns", []),
        table_data.get("rows", [])
    )


def render_flowchart(slide, data: dict):
    """フローチャートを描画"""
    steps = data.get("steps", [])
    if not steps:
        return

    header = data.get("header", {})
    step_texts = [step.get("text", step) if isinstance(step, dict) else step
                  for step in steps]

    add_flowchart(
        slide,
        step_texts,
        header.get("x_mm", 20),
        header.get("y_mm", 50) + 12,
        header.get("width_mm", 180),
        data.get("direction", "h")
    )


def render_kpi_box(slide, data: dict):
    """KPIボックスを描画"""
    value_data = data.get("value", {})

    # 背景ボックス
    add_box_with_background(
        slide,
        value_data.get("x_mm", 20) - 5,
        value_data.get("y_mm", 50) - 2,
        value_data.get("width_mm", 80) + 10,
        value_data.get("height_mm", 20) + 15
    )

    # 数値
    add_text_box(
        slide,
        value_data.get("text", ""),
        value_data.get("x_mm", 20),
        value_data.get("y_mm", 50),
        value_data.get("width_mm", 80),
        value_data.get("height_mm", 20),
        font_size_pt=value_data.get("font_size_pt", 32),
        font_bold=True,
        font_color=Colors.PRIMARY,
        alignment=PP_ALIGN.CENTER
    )

    # 単位・ラベル
    unit = data.get("unit", "")
    label = data.get("label", "")
    if unit or label:
        add_text_box(
            slide,
            f"{unit} {label}".strip(),
            value_data.get("x_mm", 20),
            value_data.get("y_mm", 50) + 22,
            value_data.get("width_mm", 80),
            10,
            font_size_pt=12,
            alignment=PP_ALIGN.CENTER
        )


def render_text_block(slide, data: dict):
    """テキストブロックを描画"""
    body = data.get("body", {})
    if body.get("text"):
        add_text_box_with_auto_shrink(
            slide,
            body["text"],
            body.get("x_mm", 20),
            body.get("y_mm", 50),
            body.get("width_mm", 180),
            body.get("height_mm", 50)
        )


def render_footer(slide, footer_data: dict):
    """フッターを描画"""
    if footer_data.get("text"):
        add_text_box(
            slide,
            footer_data["text"],
            footer_data.get("x_mm", 15),
            footer_data.get("y_mm", 275),
            footer_data.get("width_mm", 390),
            footer_data.get("height_mm", 12),
            font_size_pt=footer_data.get("font_size_pt", Fonts.FOOTER),
            font_color=Colors.SECONDARY
        )


def build_pptx(layout_data: SlideLayoutData) -> BytesIO:
    """
    LayoutEngineの出力からPPTXファイルを生成
    
    Args:
        layout_data: 座標付きレイアウトデータ
        
    Returns:
        生成されたPPTXファイルのバイナリストリーム
    """
    # プレゼンテーション作成
    prs = Presentation()
    
    # スライドサイズ設定（A3横）
    # python-pptxのデフォルト定数が機能しない場合があるため直接指定
    prs.slide_width = Mm(420)
    prs.slide_height = Mm(297)

    # 空白のレイアウトを使用（インデックス6は通常空白スライド）
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # 背景色設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = Colors.BACKGROUND

    # ヘッダー描画
    if "header" in layout_data:
        render_header(slide, layout_data["header"])

    # セクション描画
    if "sections" in layout_data:
        for section in layout_data["sections"]:
            render_section(slide, section)

    # フッター描画
    if "footer" in layout_data:
        render_footer(slide, layout_data["footer"])

    # バイトストリームに保存
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output
