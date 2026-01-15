"""
PPTX Builder Unit Tests
"""

import pytest
import sys
from pathlib import Path
from io import BytesIO

sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from pptx import Presentation
from pptx.util import Mm

from pptx_builder import (
    build_pptx,
    add_text_box,
    add_section_header,
    Colors,
    Fonts
)


class TestBuildPptx:
    """PPTX生成のテスト"""

    def test_returns_bytes(self, sample_layout_data):
        """バイトストリームを返すかテスト"""
        result = build_pptx(sample_layout_data)
        assert isinstance(result, BytesIO)
        assert result.getbuffer().nbytes > 0

    def test_valid_pptx(self, sample_layout_data):
        """有効なPPTXファイルかテスト"""
        result = build_pptx(sample_layout_data)
        # BytesIOを直接渡す
        prs = Presentation(result)
        assert prs.slide_width == Mm(420)
        assert prs.slide_height == Mm(297)

    def test_slide_size(self, sample_layout_data):
        """スライドサイズがA3横かテスト"""
        result = build_pptx(sample_layout_data)
        prs = Presentation(result)
        # A3サイズ: 420mm x 297mm
        assert int(prs.slide_width.mm) == 420
        assert int(prs.slide_height.mm) == 297

    def test_has_shapes(self, sample_layout_data):
        """図形が描画されているかテスト"""
        result = build_pptx(sample_layout_data)
        prs = Presentation(result)
        slide = prs.slides[0]
        # ヘッダー、セクション、フッターなどが描画されているはず
        assert len(slide.shapes) > 0


class TestAddTextBox:
    """テキストボックス追加のテスト"""

    @pytest.fixture
    def blank_slide(self):
        """空白スライド"""
        prs = Presentation()
        prs.slide_width = Mm(420)
        prs.slide_height = Mm(297)
        return prs.slides.add_slide(prs.slide_layouts[6])

    def test_adds_shape(self, blank_slide):
        """図形が追加される"""
        initial_count = len(blank_slide.shapes)
        add_text_box(blank_slide, "Test", 10, 10, 100, 20)
        assert len(blank_slide.shapes) == initial_count + 1

    def test_text_content(self, blank_slide):
        """テキストが正しい"""
        shape = add_text_box(blank_slide, "Hello", 10, 10, 100, 20)
        assert shape.text_frame.paragraphs[0].text == "Hello"


class TestColors:
    """カラー定義のテスト"""

    def test_primary_color(self):
        """プライマリカラーの確認"""
        assert Colors.PRIMARY is not None

    def test_from_hex(self):
        """16進数変換"""
        color = Colors.from_hex("#2B6CB0")
        assert color is not None


class TestFonts:
    """フォント定義のテスト"""

    def test_font_family(self):
        """フォントファミリーの確認"""
        assert Fonts.FAMILY == "メイリオ"

    def test_font_sizes(self):
        """フォントサイズの確認"""
        assert Fonts.TITLE == 28
        assert Fonts.SECTION_HEADER == 18
        assert Fonts.BODY == 14
